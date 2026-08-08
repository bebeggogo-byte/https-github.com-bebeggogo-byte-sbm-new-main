"""
슬라이드 ↔ 강사 발화 '깊이 차이' 분석
========================================
'보여진 슬라이드'와 '강사가 실제 말한 내용(전사본)'의 깊이 차이를 분석합니다.

찾아내는 것:
  • 말로는 풍부했지만 슬라이드엔 빈약한 부분 → 슬라이드에 채울 거리(보강 후보)
  • 슬라이드엔 있지만 거의 언급 안 한 부분 → 다음엔 더 설명하거나 덜어낼 후보
  • 즉흥적으로 더한 핵심 메시지(슬라이드에 없는 좋은 말) → 다음 버전에 반영

슬라이드 텍스트는 PPTX(텍스트박스) 또는 PPTX 내장 이미지(OCR, 선택)에서 추출.
ANTHROPIC_API_KEY 가 있으면 Claude 가 깊이 차이를 정밀 분석하고,
없으면 키워드·분량 기반 휴리스틱으로 후보를 제시합니다.

결과는 HTML 보고서(report 와 동일 테마)로 저장.

사용법:
  python gap_analysis.py --slides deck.pptx --transcript transcript.txt \
      --event output/events.json --index 0 -o output/gap_00.html
"""
from __future__ import annotations

import argparse
import html
import json
import os
import re
from datetime import datetime

import theme as theme_mod


# ---------------------------------------------------------------- 슬라이드 텍스트
def slide_texts(path: str) -> list[str]:
    """PPTX 각 슬라이드의 텍스트(텍스트박스). 이미지 덱이면 빈 문자열."""
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text.strip())
        out.append("\n".join(parts))
    return out


def ocr_slides(path: str) -> list[str]:
    """이미지 레이어 덱: 내장 이미지에서 OCR(pytesseract 필요)."""
    try:
        import pytesseract, cv2, numpy as np, zipfile
        from PIL import Image
        import io
    except ImportError:
        return []
    texts = []
    with zipfile.ZipFile(path) as z:
        media = sorted(n for n in z.namelist()
                       if n.startswith("ppt/media/") and n.lower().endswith((".png", ".jpg", ".jpeg")))
        for m in media:
            try:
                img = Image.open(io.BytesIO(z.read(m)))
                texts.append(pytesseract.image_to_string(img, lang="kor+eng").strip())
            except Exception:
                texts.append("")
    return texts


# ---------------------------------------------------------------- 분석
def analyze_ai(slides: list[str], transcript: str, meta: dict) -> dict:
    import anthropic
    client = anthropic.Anthropic()
    deck = "\n\n".join(f"[슬라이드{i+1}]\n{t or '(텍스트 없음/이미지)'}" for i, t in enumerate(slides))
    schema = (
        "강의 슬라이드와 실제 강의 전사본을 비교해 '깊이 차이'를 JSON 으로 분석:\n"
        '{"summary":[2~4줄 총평],'
        '"slide_thin_talk_rich":[{"topic":주제,"slide":슬라이드에 적힌 정도,'
        '"spoken":말로 풀어낸 깊이,"suggestion":슬라이드에 추가하면 좋을 구체적 문구/불릿}],'
        '"slide_has_talk_skipped":[{"topic":주제,"note":슬라이드엔 있으나 거의 말 안 함 → 보완/삭제 제안}],'
        '"improvised_gold":[강사가 즉흥적으로 한 좋은 말 중 다음 슬라이드에 넣을 것],'
        '"next_version_actions":[다음 버전 슬라이드 개선 액션]}'
    )
    msg = client.messages.create(
        model="claude-opus-4-8", max_tokens=4000,
        system="너는 강의 코치다. 슬라이드와 실제 발화의 깊이 차이를 짚어 슬라이드를 채우도록 돕는다. 한국어, JSON만.",
        messages=[{"role": "user", "content":
                   f"강의:{meta.get('title','')}\n{schema}\n\n=== 슬라이드 ===\n{deck}\n\n=== 전사본 ===\n{transcript[:100000]}"}],
    )
    txt = re.sub(r"```json|```", "", msg.content[0].text).strip()
    return json.loads(txt)


def analyze_heuristic(slides: list[str], transcript: str, meta: dict) -> dict:
    """API 없이: 슬라이드 키워드가 전사에서 얼마나 많이 다뤄졌는지로 갭 추정."""
    tsents = [s.strip() for s in re.split(r"(?<=[.!?。])\s+|\n+", transcript) if s.strip()]
    tlower = transcript.lower()
    thin_rich, has_skipped = [], []
    for i, st in enumerate(slides):
        # 슬라이드 핵심 단어(2글자+ 한글/영문)
        words = [w for w in re.findall(r"[가-힣A-Za-z]{2,}", st) if len(w) >= 2]
        if not words:
            continue
        spoken_hits = sum(tlower.count(w.lower()) for w in set(words))
        slide_len = len(st)
        # 슬라이드는 짧은데 관련 발화가 많음 → 채울 거리
        if slide_len < 120 and spoken_hits >= 8:
            rel = [s for s in tsents if any(w in s for w in words)][:3]
            thin_rich.append({"topic": f"슬라이드 {i+1}",
                              "slide": st[:60] + ("…" if len(st) > 60 else ""),
                              "spoken": f"전사에서 관련 언급 {spoken_hits}회",
                              "suggestion": " / ".join(rel) or "관련 발화를 불릿으로 정리"})
        # 슬라이드엔 있는데 거의 언급 없음
        if words and spoken_hits <= 1:
            has_skipped.append({"topic": f"슬라이드 {i+1}",
                               "note": f"'{words[0]}' 등 거의 언급 안 됨 → 설명 보강 또는 슬라이드 간소화"})
    return {
        "summary": [f"슬라이드 {len(slides)}장 · 전사 {len(tsents)}문장 비교(휴리스틱).",
                    "ANTHROPIC_API_KEY 설정 시 Claude 가 깊이 차이를 정밀 분석합니다."],
        "slide_thin_talk_rich": thin_rich[:8],
        "slide_has_talk_skipped": has_skipped[:6],
        "improvised_gold": [],
        "next_version_actions": ["풍부히 말한 부분을 슬라이드 불릿으로 승격",
                                 "언급 적은 슬라이드는 설명 보강 또는 삭제"],
    }


# ---------------------------------------------------------------- HTML
def _li(items): return "\n".join(f"<li>{html.escape(str(x))}</li>" for x in items)


def render(data: dict, meta: dict, th: dict) -> str:
    pal, fonts = th["palette"], th["fonts"]
    def block_tr(rows):
        return "\n".join(
            f'<div class="card"><b>{html.escape(r.get("topic",""))}</b>'
            f'<div class="row"><span class="tag slide">슬라이드</span> {html.escape(r.get("slide",""))}</div>'
            f'<div class="row"><span class="tag talk">발화</span> {html.escape(r.get("spoken",""))}</div>'
            f'<div class="sug">➕ {html.escape(r.get("suggestion",""))}</div></div>' for r in rows)
    skipped = "\n".join(f'<div class="card"><b>{html.escape(r.get("topic",""))}</b>'
                        f'<div class="sug">{html.escape(r.get("note",""))}</div></div>'
                        for r in data.get("slide_has_talk_skipped", []))
    return f"""<!doctype html><html lang="ko"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html.escape(meta.get('title','강의'))} — 깊이 차이 분석</title>
<style>
 body{{margin:0;background:#{pal['bg']};color:#{pal['ink']};
   font-family:'{fonts['body']}','Apple SD Gothic Neo',sans-serif;line-height:1.7}}
 .wrap{{max-width:880px;margin:0 auto;padding:48px 28px 80px}}
 h1{{font-family:'{fonts['heading']}',serif;color:#{pal['heading']};font-size:27px;margin:0 0 4px}}
 .meta{{color:#{pal['muted']};font-size:14px;margin-bottom:18px}}
 h2{{font-family:'{fonts['heading']}',serif;color:#{pal['heading']};font-size:19px;margin:34px 0 10px}}
 .tldr{{background:color-mix(in srgb,#{pal['accent']} 8%,#{pal['bg']});
   border-left:4px solid #{pal['accent']};border-radius:8px;padding:14px 18px}}
 .card{{border:1px solid #{pal['line']};border-radius:10px;padding:14px 16px;margin:10px 0}}
 .row{{margin:4px 0}} .tag{{display:inline-block;font-size:11px;padding:1px 8px;border-radius:10px;
   color:#fff;margin-right:6px}} .tag.slide{{background:#{pal['muted']}}} .tag.talk{{background:#{pal['accent']}}}
 .sug{{margin-top:6px;color:#{pal['heading']};font-weight:600}}
 ul{{padding-left:20px}} footer{{margin-top:40px;text-align:right;color:#{pal['muted']};
   font-size:11px;border-top:1px solid #{pal['line']};padding-top:10px}}
</style></head><body><div class="wrap">
 <h1>{html.escape(meta.get('title','강의'))} — 슬라이드 vs 강의 깊이 차이</h1>
 <div class="meta">{html.escape(meta.get('date',''))} · 생성 {datetime.now():%Y-%m-%d}</div>
 <h2>총평</h2><div class="tldr"><ul>{_li(data.get('summary',[]))}</ul></div>
 <h2>🗣️ 말로는 깊었지만 슬라이드엔 빈약 → 채울 거리</h2>
 {block_tr(data.get('slide_thin_talk_rich',[])) or '<p>해당 없음</p>'}
 <h2>📄 슬라이드엔 있으나 거의 말하지 않음</h2>
 {skipped or '<p>해당 없음</p>'}
 <h2>✨ 즉흥적으로 나온 좋은 말 (다음 슬라이드에 반영)</h2>
 <ul>{_li(data.get('improvised_gold',[])) or '<li>(전사 분석에서 추출되면 표시)</li>'}</ul>
 <h2>✅ 다음 버전 개선 액션</h2><ul>{_li(data.get('next_version_actions',[]))}</ul>
 <footer>{html.escape(th.get('source_label',''))}</footer>
</div></body></html>"""


def main() -> None:
    ap = argparse.ArgumentParser(description="슬라이드 vs 발화 깊이 차이 분석")
    ap.add_argument("--slides", required=True, help="강의에 쓴 PPTX")
    ap.add_argument("--transcript", required=True, help="강의 전사본 .txt")
    ap.add_argument("--event"); ap.add_argument("--index", type=int, default=0)
    ap.add_argument("--ocr", action="store_true", help="이미지 덱이면 OCR로 텍스트 추출")
    ap.add_argument("--theme")
    ap.add_argument("-o", "--output", default="output/gap.html")
    args = ap.parse_args()

    slides = slide_texts(args.slides)
    if args.ocr and sum(len(s) for s in slides) < 50:
        slides = ocr_slides(args.slides) or slides
    transcript = open(args.transcript, encoding="utf-8").read()

    meta = {"title": "강의"}
    if args.event:
        meta = json.load(open(args.event, encoding="utf-8"))[args.index]
    th = theme_mod.resolve(meta, forced=args.theme)
    print(f"🎨 테마: {th['display_name']}")

    if os.getenv("ANTHROPIC_API_KEY"):
        try:
            data = analyze_ai(slides, transcript, meta); print("🧠 Claude 깊이 분석 완료")
        except Exception as e:
            print(f"⚠️ AI 분석 실패({e}) → 휴리스틱"); data = analyze_heuristic(slides, transcript, meta)
    else:
        print("⚠️ ANTHROPIC_API_KEY 미설정 → 휴리스틱 분석")
        data = analyze_heuristic(slides, transcript, meta)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    open(args.output, "w", encoding="utf-8").write(render(data, meta, th))
    print(f"✅ 깊이 차이 분석 → {args.output}")


if __name__ == "__main__":
    main()
