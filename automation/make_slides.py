"""
커리큘럼 → 강의 슬라이드(PPTX) 생성  [고객사 테마 적용]
=======================================================
레이아웃 사양(고정):
  • 배경: 흰색
  • 제목 / 소제목 / 학습목표: 왼쪽 상단
  • 본문(핵심내용): 그 아래
  • 출처: 오른쪽 하단, 폰트 6pt (텍스트는 고객사 테마의 source_label)

색·폰트·출처라벨은 고객사 테마(themes/<name>.json)에서 옵니다.
테마는 --theme 로 강제 지정하거나, 미지정 시 일정 내용으로 자동 분류됩니다.

사용법:
    python make_slides.py --from-events output/events.json --index 0 -o output/slides.pptx
    python make_slides.py --from-events output/events.json --index 0 --theme nedabah -o out.pptx
    python make_slides.py --slides output/slides.json --theme acme -o output/slides.pptx
"""
from __future__ import annotations

import argparse
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

import theme as theme_mod


def C(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _txt(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _run(p, text, size, color, font, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return r


def add_slide(prs, s: dict, th: dict) -> None:
    pal = th["palette"]
    fonts = th["fonts"]
    lay = th["layout"]
    heading_f, body_f = fonts["heading"], fonts["body"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(pal["bg"])

    # ── 왼쪽 상단: 제목 / 소제목 / 학습목표 ──
    tf = _txt(slide, Inches(0.55), Inches(0.45), Inches(8.4), Inches(2.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, s.get("title", ""), lay["title_pt"], C(pal["heading"]), heading_f, bold=True)

    if s.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(2)
        _run(p2, s["subtitle"], lay["subtitle_pt"], C(pal["accent"]), body_f)

    if s.get("objective"):
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        p3.space_before = Pt(6)
        _run(p3, f"학습목표  {s['objective']}", lay["objective_pt"], C(pal["muted"]), body_f)

    # ── 본문(핵심내용) ──
    if s.get("bullets"):
        body = _txt(slide, Inches(0.6), Inches(2.7), Inches(8.8), Inches(4.2))
        for i, b in enumerate(s["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
            _run(p, "•  ", lay["body_pt"], C(pal["accent"]), body_f, bold=True)
            _run(p, b, lay["body_pt"], C(pal["ink"]), body_f)

    # ── 오른쪽 하단: 출처(테마 라벨, 6pt) ──
    label = th.get("source_label", "")
    if label:
        src = _txt(slide, Inches(6.3), Inches(7.05), Inches(3.3), Inches(0.35),
                   anchor=MSO_ANCHOR.BOTTOM)
        p = src.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _run(p, label, lay["source_pt"], C(pal["muted"]), body_f)


def slides_from_event(e: dict) -> list[dict]:
    """events.json 1건 → 대상(공공기관/학생/일반) 맞춤 슬라이드 개요 자동 생성."""
    import audience as aud
    title = e.get("title", "강의")
    subject = e.get("subject") or title
    who = e.get("audience") or "참석자"
    name, p = aud.resolve(e)
    struct = p["structure"]  # 대상별 권장 구성(섹션 제목)

    # 표지 + 구성 단계별 슬라이드
    slides = [{"title": title, "subtitle": subject,
               "objective": f"{p['display']} 대상 · {p['objective_style']}",
               "bullets": [f"오늘 함께 볼 것: {' → '.join(struct)}",
                           "왜 지금 이 주제인가",
                           f"진행 방식: {p['interaction']} 상호작용"]}]
    for sec in struct:
        slides.append({
            "title": sec, "subtitle": subject,
            "objective": p["objective_style"],
            "bullets": [f"{sec} 핵심 1", f"{sec} 핵심 2",
                        f"{who} 맞춤: {p['examples']} 활용"]})
    return slides


def main() -> None:
    ap = argparse.ArgumentParser(description="커리큘럼 → 슬라이드 PPTX (고객사 테마)")
    ap.add_argument("--slides", help="슬라이드 JSON 파일")
    ap.add_argument("--from-events", help="events.json 에서 자동 생성")
    ap.add_argument("--index", type=int, default=0)
    ap.add_argument("--theme", help="고객사 테마 강제 지정(미지정 시 자동 분류)")
    ap.add_argument("-o", "--output", default="output/slides.pptx")
    args = ap.parse_args()

    event = {}
    if args.slides:
        slides = json.load(open(args.slides, encoding="utf-8"))
    elif args.from_events:
        events = json.load(open(args.from_events, encoding="utf-8"))
        event = events[args.index]
        slides = slides_from_event(event)
    else:
        ap.error("--slides 또는 --from-events 중 하나가 필요합니다.")

    th = theme_mod.resolve(event, forced=args.theme)
    print(f"🎨 테마: {th['display_name']} — {th.get('_resolved_by','')}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for s in slides:
        add_slide(prs, s, th)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    print(f"✅ 슬라이드 {len(slides)}장 → {args.output}")


if __name__ == "__main__":
    main()
