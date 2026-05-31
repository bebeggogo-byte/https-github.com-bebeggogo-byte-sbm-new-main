"""
커리큘럼 → 강의 워크지(워크시트) 생성  [고객사 테마 적용]
=========================================================
참석자가 작성하는 워크지를 PPTX(편집용) + PDF(인쇄용) 로 만듭니다.
색·출처라벨은 고객사 테마(themes/<name>.json)에서 오며, --theme 강제 지정 또는 자동 분류.

사용법:
    python make_worksheet.py --from-events output/events.json --index 0 \
        -o output/worksheet.pptx --pdf output/worksheet.pdf [--theme nedabah]
"""
from __future__ import annotations

import argparse
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

import theme as theme_mod


def C(h: str) -> RGBColor:
    return RGBColor.from_string(h)


def _hex_to_rgb01(h: str) -> tuple[float, float, float]:
    return tuple(int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))


def worksheet_items(e: dict) -> dict:
    subject = e.get("subject") or e.get("title", "강의")
    audience = e.get("audience") or "참석자"
    return {
        "title": f"{e.get('title','강의')} 워크지",
        "meta": " · ".join(x for x in [e.get("date", ""), e.get("start_time", ""),
                                       f"대상 {audience}"] if x),
        "sections": [
            ("들어가며", "오늘 이 주제에 대해 내가 가진 질문이나 기대를 적어보세요."),
            ("핵심 개념 정리", f"'{subject}'의 핵심 개념을 내 언어로 한 문장으로 요약하면?"),
            ("근거 / 출처", "오늘 들은 내용 중 가장 마음에 남는 부분과 이유는?"),
            ("나의 적용", "내 삶/업무의 어느 영역에 어떻게 적용할 수 있을까요? (구체적으로)"),
            ("팀·공동체 적용", "함께 실천할 수 있는 한 가지는?"),
            ("한 줄 결단", "오늘 강의 후 나의 결단을 한 문장으로 적어보세요."),
        ],
    }


def build_pptx(data: dict, out: str, th: dict) -> None:
    pal, fonts = th["palette"], th["fonts"]
    heading_f, body_f = fonts["heading"], fonts["body"]

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(pal["bg"])

    head = slide.shapes.add_textbox(Inches(0.55), Inches(0.4), Inches(8.9), Inches(1.0))
    tf = head.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = data["title"]
    r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = C(pal["heading"]); r.font.name = heading_f
    if data["meta"]:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = data["meta"]
        r2.font.size = Pt(11); r2.font.color.rgb = C(pal["muted"]); r2.font.name = body_f

    top = 1.55
    for i, (head_t, q) in enumerate(data["sections"]):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.85))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{i+1}. {head_t}"
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = C(pal["ink"]); r.font.name = body_f
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = q
        r2.font.size = Pt(11); r2.font.color.rgb = C(pal["muted"]); r2.font.name = body_f
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6),
                                    Inches(top + 0.78), Inches(8.8), Pt(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = C(pal["line"]); ln.line.fill.background()
        top += 0.92

    label = th.get("source_label", "")
    if label:
        box = slide.shapes.add_textbox(Inches(6.3), Inches(7.05), Inches(3.3), Inches(0.35))
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = label
        r.font.size = Pt(th["layout"]["source_pt"])
        r.font.color.rgb = C(pal["muted"]); r.font.name = body_f

    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    prs.save(out)
    print(f"✅ 워크지(PPTX) → {out}")


def build_pdf(data: dict, out: str, th: dict) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfgen import canvas

    pal = th["palette"]
    f = "HYSMyeongJo-Medium"  # 한글 CID 폰트(폰트파일 불필요)
    pdfmetrics.registerFont(UnicodeCIDFont(f))
    c = canvas.Canvas(out, pagesize=A4)
    w, h = A4
    y = h - 25 * mm
    c.setFont(f, 18); c.setFillColorRGB(*_hex_to_rgb01(pal["heading"]))
    c.drawString(22 * mm, y, data["title"])
    y -= 8 * mm
    c.setFont(f, 9); c.setFillColorRGB(*_hex_to_rgb01(pal["muted"]))
    c.drawString(22 * mm, y, data["meta"])
    y -= 12 * mm
    for i, (head_t, q) in enumerate(data["sections"]):
        c.setFont(f, 12); c.setFillColorRGB(*_hex_to_rgb01(pal["ink"]))
        c.drawString(22 * mm, y, f"{i+1}. {head_t}")
        y -= 6 * mm
        c.setFont(f, 9.5); c.setFillColorRGB(*_hex_to_rgb01(pal["muted"]))
        c.drawString(24 * mm, y, q)
        y -= 9 * mm
        c.setStrokeColorRGB(*_hex_to_rgb01(pal["line"]))
        for _ in range(2):
            c.line(24 * mm, y, w - 22 * mm, y)
            y -= 8 * mm
        y -= 4 * mm
    label = th.get("source_label", "")
    if label:
        c.setFont(f, th["layout"]["source_pt"])
        c.setFillColorRGB(*_hex_to_rgb01(pal["muted"]))
        c.drawRightString(w - 22 * mm, 14 * mm, label)
    c.save()
    print(f"✅ 워크지(PDF) → {out}")


def main() -> None:
    ap = argparse.ArgumentParser(description="커리큘럼 → 워크지 (고객사 테마)")
    ap.add_argument("--from-events", required=True)
    ap.add_argument("--index", type=int, default=0)
    ap.add_argument("--theme", help="고객사 테마 강제 지정(미지정 시 자동 분류)")
    ap.add_argument("-o", "--output", default="output/worksheet.pptx")
    ap.add_argument("--pdf", help="PDF 도 함께 생성")
    args = ap.parse_args()

    events = json.load(open(args.from_events, encoding="utf-8"))
    e = events[args.index]
    th = theme_mod.resolve(e, forced=args.theme)
    print(f"🎨 테마: {th['display_name']} — {th.get('_resolved_by','')}")
    data = worksheet_items(e)
    build_pptx(data, args.output, th)
    if args.pdf:
        build_pdf(data, args.pdf, th)


if __name__ == "__main__":
    main()
